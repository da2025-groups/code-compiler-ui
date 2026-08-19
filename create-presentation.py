#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Code Compiler Platform
Uses python-pptx for proper slide layout without cropping issues
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (Indigo)
INDIGO = RGBColor(79, 70, 229)  # #4f46e5
DARK_INDIGO = RGBColor(55, 48, 163)  # #3730a3
LIGHT_PURPLE = RGBColor(237, 233, 254)  # #ede9fe
DARK_TEXT = RGBColor(17, 24, 39)  # #111827
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = INDIGO

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = RGBColor(199, 210, 254)  # Light indigo

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Team: Bala Vardhan Palli and Team\nProject #10 — Dr. Dhawaleswar Rao, SoET/CSE"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = RGBColor(199, 210, 254)

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = INDIGO

    # Content area
    content_top = Inches(1.4)

    for item in content_items:
        if item['type'] == 'text':
            text_box = slide.shapes.add_textbox(
                Inches(item.get('left', 0.8)),
                content_top,
                Inches(item.get('width', 8.4)),
                Inches(item.get('height', 0.6))
            )
            text_frame = text_box.text_frame
            text_frame.text = item['text']
            para = text_frame.paragraphs[0]
            para.font.size = Pt(item.get('size', 18))
            para.font.bold = item.get('bold', False)
            para.font.color.rgb = item.get('color', DARK_TEXT)
            content_top += Inches(item.get('height', 0.6) + 0.1)

        elif item['type'] == 'bullets':
            for bullet in item['bullets']:
                bullet_box = slide.shapes.add_textbox(
                    Inches(1.2),
                    content_top,
                    Inches(8),
                    Inches(0.4)
                )
                bullet_frame = bullet_box.text_frame
                bullet_frame.text = f"• {bullet}"
                bullet_para = bullet_frame.paragraphs[0]
                bullet_para.font.size = Pt(16)
                bullet_para.font.color.rgb = DARK_TEXT
                content_top += Inches(0.5)

        elif item['type'] == 'box':
            # Colored box
            box_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(item.get('left', 0.8)),
                content_top,
                Inches(item.get('width', 8.4)),
                Inches(item.get('height', 1.2))
            )
            fill = box_shape.fill
            fill.solid()
            fill.fore_color.rgb = item.get('bg_color', LIGHT_PURPLE)

            # Box text
            text_frame = box_shape.text_frame
            text_frame.text = item['text']
            para = text_frame.paragraphs[0]
            para.font.size = Pt(item.get('size', 16))
            para.font.bold = item.get('bold', False)
            para.font.color.rgb = item.get('color', DARK_TEXT)

            content_top += Inches(item.get('height', 1.2) + 0.2)

    return slide

# Slide 1: Cover
add_title_slide(
    prs,
    "Code Compiler Platform",
    "HackerRank-style Competitive Coding System\nAutomated Code Evaluation · Multi-language Support · Live Leaderboard"
)

# Slide 2: System Architecture
arch_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = arch_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "System Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = INDIGO

y = Inches(1.2)

# Frontend box
frontend_box = arch_slide.shapes.add_shape(1, Inches(0.8), y, Inches(8.4), Inches(1.3))
fill = frontend_box.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PURPLE
text_frame = frontend_box.text_frame
text_frame.text = "🎨 FRONTEND — React + Vite + Monaco Editor\n\nStudent: Playground · Problems · Code Editor · Leaderboard\nAdmin: Question CRUD · Test Cases · Submissions Monitor"
para = text_frame.paragraphs[0]
para.font.size = Pt(14)
para.font.bold = True
para.font.color.rgb = DARK_INDIGO

y += Inches(1.5)

# Arrow
arrow_box = arch_slide.shapes.add_textbox(Inches(4), y, Inches(2), Inches(0.4))
arrow_frame = arrow_box.text_frame
arrow_frame.text = "⬇️ REST API + JWT ⬇️"
arrow_para = arrow_frame.paragraphs[0]
arrow_para.alignment = PP_ALIGN.CENTER
arrow_para.font.size = Pt(16)
arrow_para.font.bold = True
arrow_para.font.color.rgb = INDIGO

y += Inches(0.5)

# Backend box
backend_box = arch_slide.shapes.add_shape(1, Inches(0.8), y, Inches(8.4), Inches(1.3))
fill = backend_box.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PURPLE
text_frame = backend_box.text_frame
text_frame.text = "⚙️ BACKEND — FastAPI + SQLite\n\nRoutes: /auth · /questions · /submissions · /rankings · /admin\nScoring: (passed_cases / total_cases) × 100 · DB: users, questions, submissions"
para = text_frame.paragraphs[0]
para.font.size = Pt(14)
para.font.bold = True
para.font.color.rgb = DARK_INDIGO

y += Inches(1.5)

# Arrow
arrow_box2 = arch_slide.shapes.add_textbox(Inches(4), y, Inches(2), Inches(0.4))
arrow_frame2 = arrow_box2.text_frame
arrow_frame2.text = "⬇️ HTTP (Docker) ⬇️"
arrow_para2 = arrow_frame2.paragraphs[0]
arrow_para2.alignment = PP_ALIGN.CENTER
arrow_para2.font.size = Pt(16)
arrow_para2.font.bold = True
arrow_para2.font.color.rgb = INDIGO

y += Inches(0.5)

# Piston box
piston_box = arch_slide.shapes.add_shape(1, Inches(0.8), y, Inches(8.4), Inches(1.3))
fill = piston_box.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PURPLE
text_frame = piston_box.text_frame
text_frame.text = "🔧 PISTON ENGINE — Multi-language Sandbox\n\nInput: { language, code, stdin } → Output: { stdout, stderr, time }\nLanguages: Python · JavaScript · Java · C++ · Go · Rust"
para = text_frame.paragraphs[0]
para.font.size = Pt(14)
para.font.bold = True
para.font.color.rgb = DARK_INDIGO

# Slide 3: Tech Stack & Demo
tech_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = tech_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Tech Stack & Quick Start"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = INDIGO

y = Inches(1.1)

# Tech stack
tech_box = tech_slide.shapes.add_textbox(Inches(0.8), y, Inches(4), Inches(2.5))
tech_frame = tech_box.text_frame
tech_frame.text = """Frontend: React 18 · Vite · Material UI · Monaco Editor
Backend: FastAPI · SQLite · SQLAlchemy · JWT
Execution: Piston (Docker) — 6 languages
Deployment: Docker Compose · Makefile"""
para = tech_frame.paragraphs[0]
para.font.size = Pt(14)
para.line_spacing = 1.5
para.font.color.rgb = DARK_TEXT

# Quick start
start_box = tech_slide.shapes.add_textbox(Inches(5.2), y, Inches(4), Inches(2.5))
start_frame = start_box.text_frame
start_frame.text = """Quick Start:
  make setup        # Start backend
  make seed         # Create questions
  make frontend-dev # Launch UI

Demo: http://localhost:5173
Login: admin@platform.com / admin123"""
para = start_frame.paragraphs[0]
para.font.size = Pt(14)
para.line_spacing = 1.5
para.font.color.rgb = DARK_TEXT

y += Inches(2.7)

# Features
features_box = tech_slide.shapes.add_shape(1, Inches(0.8), y, Inches(8.4), Inches(2))
fill = features_box.fill
fill.solid()
fill.fore_color.rgb = LIGHT_PURPLE
text_frame = features_box.text_frame
text_frame.text = """✅ Playground — Free editor with custom stdin
✅ Questions — 11 problems (Easy/Medium) with sample I/O
✅ Live Judge — Instant verdict table per test case
✅ Leaderboard — Ranked by total score, then solved count
✅ Admin Panel — Create/edit questions with JSON test cases"""
para = text_frame.paragraphs[0]
para.font.size = Pt(16)
para.line_spacing = 1.3
para.font.color.rgb = DARK_TEXT

# Save
prs.save('docs/presentation-python.pptx')
print("✅ Created docs/presentation-python.pptx")
